import logging
import json
import os
from pathlib import Path
import warnings
import torch
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
import hashlib
import re
# from bs4 import BeautifulSoup
import html2text
from PIL import Image, ImageStat
import io
import numpy as np
import cv2
import math


os.environ["HF_HUB_DISABLE_HARDLINKS"] = "1"
os.environ["HF_HUB_DISABLE_SYMLINKS"] = "1" 

# Ensure PyTorch uses GPU if available
torch.set_default_device("cuda" if torch.cuda.is_available() else "cpu")

print("CUDA Available:", torch.cuda.is_available())
print("Using Device:", torch.cuda.get_device_name(0) if torch.cuda.is_available() else "CPU")


from docling.datamodel.base_models import InputFormat
from docling_core.types.doc import ImageRefMode
from docling.document_converter import DocumentConverter, PdfFormatOption, WordFormatOption
from docling.pipeline.standard_pdf_pipeline import StandardPdfPipeline
from docling.datamodel.pipeline_options import PdfPipelineOptions
from docling.pipeline.simple_pipeline import SimplePipeline
from docling.datamodel.settings import settings

warnings.filterwarnings("ignore")

# Logging setup
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

IMAGE_RESOLUTION_SCALE = 2.0

def create_pipeline_options(input_format):
    if input_format == InputFormat.PDF:
        return PdfFormatOption(
            pipeline_options=PdfPipelineOptions(
                do_table_structure=True,
                generate_page_images=True,
                generate_picture_images=True,
                images_scale=IMAGE_RESOLUTION_SCALE,
            )
        )
    elif input_format == InputFormat.DOCX:
        return WordFormatOption(pipeline_cls=SimplePipeline)
    return None


def initialize_converter():
    allowed_formats = [InputFormat.PDF, InputFormat.DOCX]
    format_options = {
        fmt: create_pipeline_options(fmt)
        for fmt in allowed_formats if create_pipeline_options(fmt)
    }
    return DocumentConverter(allowed_formats=allowed_formats, format_options=format_options)


def convert_and_save(input_file: Path, output_dir: Path):
    """Converts documents to Markdown and saves the output."""
    os.makedirs(output_dir, exist_ok=True)
    input_paths = [input_file]
    # input_paths = list(Path(INPUT_DIR).glob("*.*"))

    if not input_paths:
        logger.warning("No input files found in the directory.")
        return {"success": False, "message": "No input files found."}

    converter = initialize_converter()
    docling_paths = [p for p in input_paths if p.suffix.lower() in [".pdf", ".docx"]]
    other_paths = [p for p in input_paths if p.suffix.lower() in [".csv", ".xlsx", ".txt", ".pptx", ".json", ".html", ".md"]]

    processed_files = []

    try:
        # Run Docling conversions
        if docling_paths:
            results = converter.convert_all(docling_paths)
            for res in results:
                file_name = res.input.file.stem
                md_path = output_dir / f"{file_name}.md"
                res.document.save_as_markdown(md_path, image_mode=ImageRefMode.REFERENCED)
                processed_files.append(md_path.name)
                logger.info(f"Markdown saved to {md_path}")

        # Run fallback conversions
        for file_path in other_paths:
            file_name = file_path.stem
            md_path = output_dir / f"{file_name}.md"
            convert_non_docling(file_path, md_path)
            processed_files.append(md_path.name)
            logger.info(f"Markdown saved to {md_path} (fallback parser)")

        if not processed_files:
            return {"success": False, "message": "No files were converted."}

        return {"success": True, "message": "Conversion completed successfully.", "files": processed_files}

    except Exception as e:
        logger.error(f"Error during conversion: {e}")
        return {"success": False, "message": f"Conversion failed: {e}"}



def convert_non_docling(file_path, output_path):
    ext = file_path.suffix.lower()

    if ext == ".csv":
        df = pd.read_csv(file_path)
        md_parts = []
        for i, row in df.iterrows():
            md_parts.append(f"## Data {i + 1}")
            for col, val in row.items():
                md_parts.append(f"- {col}: {val}")
            md_parts.append("")  # spacing
        md = "\n".join(md_parts)

    elif ext == ".xlsx":
        xls = pd.ExcelFile(file_path)
        md_parts = []
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            # md_parts.append(f"# Sheet name: {sheet_name}\n")
            for i, row in df.iterrows():
                md_parts.append(f"## Data {i + 1}")
                for col, val in row.items():
                    md_parts.append(f"- {col}: {val}")
                md_parts.append("")  # spacing
        md = "\n".join(md_parts)

    elif ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read().strip()

        paragraphs = [p.strip() for p in re.split(r'\n\s*\n', content) if p.strip()]

        md_lines = []
        for i, para in enumerate(paragraphs, start=1):
            md_lines.append(f"## Paragraph {i}")
            md_lines.append(para)
            md_lines.append("")  # blank line between sections

        md = "\n".join(md_lines)

    elif ext == ".pptx":
        md = parse_ppt(file_path, output_path.parent)

    elif ext == ".json":
        md = parse_json(file_path)

    elif ext == ".md":
        with open(file_path, "r", encoding="utf-8") as f:
            md = f.read()

    elif ext == ".html":
        md = parse_html(file_path)

    else:
        return  # unsupported

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(md)


def extract_images_from_shape(shape, prs, artifact_dir, image_counter, seen_hashes, text_shapes):
    images = []

    # ensure shape is picture
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return images

    try:
        image_bytes = shape.image.blob
        img_hash = hashlib.sha256(image_bytes).hexdigest()
        if img_hash in seen_hashes:
            return images
        seen_hashes.add(img_hash)

        # open image WITHOUT converting away alpha
        img = Image.open(io.BytesIO(image_bytes))
        width, height = img.size
        aspect_ratio = width / height if height else 1

        # --- Rule: skip extreme aspect ratios (banners/lines) ---
        if aspect_ratio > 6 or aspect_ratio < 0.16:
            return images

        # --- Rule: skip near-solid-color images ---
        # convert to RGB for color stats if needed (but keep original img for alpha)
        stat = ImageStat.Stat(img.convert("RGB"))
        stddev = sum(stat.stddev) / 3
        if stddev < 10:
            return images

        # --- Rule: skip small relative to slide size ---
        slide_w, slide_h = prs.slide_width, prs.slide_height
        rel_w = shape.width / slide_w
        rel_h = shape.height / slide_h
        # use area fraction too
        if (shape.width * shape.height) < 0.02 * (slide_w * slide_h):
            return images

        # --- Rule: skip images that are mostly transparent ---
        # check for alpha channel robustly
        alpha_ratio = 0.0
        # Pillow band names may include 'A'
        bands = img.getbands()  # e.g. ('R','G','B') or ('R','G','B','A')
        if 'A' in bands:
            alpha = np.array(img.getchannel('A'))  # 0..255, 0=transparent
            # compute proportion of pixels that are transparent (or nearly transparent)
            alpha_ratio = float(np.mean(alpha < 250))  # threshold can be tuned
            if alpha_ratio > 0.00001:  # >30% transparent -> skip
                return images
        # else: no alpha channel -> we cannot treat as transparent

        # --- Rule: skip images covered by text shapes (overlap) ---
        img_left, img_top = shape.left, shape.top
        img_right = img_left + shape.width
        img_bottom = img_top + shape.height
        img_area = shape.width * shape.height if (shape.width and shape.height) else 1

        for txt in text_shapes:
            t_left, t_top = txt.left, txt.top
            t_right = t_left + txt.width
            t_bottom = t_top + txt.height

            overlap_x = max(0, min(img_right, t_right) - max(img_left, t_left))
            overlap_y = max(0, min(img_bottom, t_bottom) - max(img_top, t_top))
            overlap_area = overlap_x * overlap_y
            if img_area > 0 and (overlap_area / img_area) > 0.4:  # >40% covered
                return images

        # --- Passed filters: save image ---
        # preserve alpha if present by saving PNG with alpha; otherwise convert to RGB
        image_ext = shape.image.ext or "png"
        img_filename = f"image_{image_counter:06d}_{img_hash[:16]}.{image_ext}"
        img_path = artifact_dir / img_filename

        # If image has alpha and ext supports it, save preserving alpha; otherwise convert
        if 'A' in bands and image_ext.lower() in ("png",):
            img.save(img_path, format="PNG")
        else:
            # convert to RGB to be safe (some formats may be JPEG)
            img.convert("RGB").save(img_path, format="PNG")

        images.append((str(img_path), image_counter + 1))

    except Exception as e:
        print(f"Error extracting image: {e}")

    return images

def parse_ppt(file_path, output_dir):
    prs = Presentation(file_path)
    artifact_dir = Path(output_dir) / f"{Path(file_path).stem}_artifacts"
    artifact_dir.mkdir(parents=True, exist_ok=True)

    slides = []
    image_counter = 0
    seen_hashes = set()

    for i, slide in enumerate(prs.slides, start=1):
        slide_content = [f"## Slide {i}"]

        # Collect text shapes for overlap detection
        text_shapes = [s for s in slide.shapes if hasattr(s, "text") and s.text.strip()]

        for shape in slide.shapes:
            # Extract text
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())

            # ✅ Pass prs as an argument here
            extracted = extract_images_from_shape(
                shape, prs, artifact_dir, image_counter, seen_hashes, text_shapes
            )

            for img_path, new_counter in extracted:
                slide_content.append(f"![Image]({img_path})")
                image_counter = new_counter

        slides.append("\n".join(slide_content))

    md = "\n\n".join(slides)
    return md


def parse_json(file_path):
    """
    Parse a JSON file into Markdown format.
    Converts keys and values into readable sections.
    Handles both dicts and lists.
    """
    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    def format_json(data, level=2):
        """Recursively format JSON as Markdown text."""
        md_lines = []
        if isinstance(data, dict):
            for key, value in data.items():
                header = "#" * level + f" {key}"
                md_lines.append(header)
                md_lines.append("")
                md_lines.append(format_json(value, level + 1))
        elif isinstance(data, list):
            for i, item in enumerate(data, start=1):
                md_lines.append(f"- **Item {i}:**")
                md_lines.append(format_json(item, level + 1))
        else:
            md_lines.append(str(data))
        return "\n".join(md_lines)

    md = format_json(data)
    return md

def parse_html(file_path):
    """
    Parse an HTML file and convert it to Markdown.
    Keeps text, headings, links, and images in readable Markdown format.
    """
    with open(file_path, "r", encoding="utf-8") as f:
        html_content = f.read()

    h = html2text.HTML2Text()
    h.ignore_links = False
    h.ignore_images = False
    h.ignore_emphasis = False
    h.bypass_tables = False
    h.body_width = 0  # Prevent wrapping

    md = h.handle(html_content).strip()

    # Optional: cleanup empty lines
    md = "\n".join([line.strip() for line in md.splitlines() if line.strip()])

    return md